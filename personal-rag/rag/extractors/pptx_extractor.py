from __future__ import annotations

from pathlib import Path
from pptx import Presentation

from rag.core.models import Page


def extract_pptx(path: Path, doc_id: str) -> list[Page]:
    presentation = Presentation(str(path))
    pages: list[Page] = []
    for idx, slide in enumerate(presentation.slides, start=1):
        text_runs: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_runs.append(shape.text)
        pages.append(Page(doc_id, str(path), path.name, idx, f"slide {idx}", "\n".join(text_runs).strip(), {"type": "pptx"}))
    return pages
